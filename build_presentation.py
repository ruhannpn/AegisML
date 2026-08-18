import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Dark Enterprise Glassmorphism Theme)
    COLOR_BG = RGBColor(9, 13, 22)         # #090d16
    COLOR_CARD = RGBColor(17, 24, 39)      # #111827
    COLOR_TEXT = RGBColor(243, 244, 246)   # #f3f4f6
    COLOR_MUTED = RGBColor(156, 163, 175)  # #9ca3af
    COLOR_PURPLE = RGBColor(124, 58, 237)  # #7c3aed
    COLOR_ACCENT = RGBColor(99, 102, 241)  # #6366f1
    COLOR_GREEN = RGBColor(16, 185, 129)   # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)   # #f59e0b
    COLOR_RED = RGBColor(239, 68, 68)      # #ef4444
    COLOR_BLUE = RGBColor(59, 130, 246)    # #3b82f6
    COLOR_CYAN = RGBColor(6, 182, 212)     # #06b6d4

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="Review 2 Parameter"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Inter"

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT
        p_title.font.name = "Inter"

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.3), Inches(10.333), Inches(4.9))
    card.fill.solid(); card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_PURPLE; card.line.width = Pt(2)

    tf1 = card.text_frame; tf1.word_wrap = True
    tf1.margin_left = Inches(0.5); tf1.margin_top = Inches(0.6)

    p0 = tf1.paragraphs[0]
    p0.text = "REVIEW 2 ACADEMIC EVALUATION PRESENTATION"
    p0.font.size = Pt(11); p0.font.bold = True; p0.font.color.rgb = COLOR_AMBER; p0.font.name = "Inter"

    p1 = tf1.add_paragraph()
    p1.text = "AegisML: Autonomous Multi-Agent AI Data Science & Governance Platform"
    p1.font.size = Pt(26); p1.font.bold = True; p1.font.color.rgb = COLOR_TEXT; p1.font.name = "Inter"

    p2 = tf1.add_paragraph()
    p2.text = "\nAn auditable, human-in-the-loop multi-agent architecture featuring 6 specialized pipeline nodes (Data Analysis Agent, Planner Agent, Data Agent, Training Agent, Fairness Agent, Governance Gate) with automated quality retry loops and interactive human governance loops."
    p2.font.size = Pt(13); p2.font.color.rgb = COLOR_MUTED; p2.font.name = "Inter"

    # -------------------------------------------------------------
    # SLIDE 2: Parameter 1 — Domain Understanding & Problem Definition
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Parameter 1: Domain Understanding & Problem Definition (3 Marks)", "CO2 Mapped")

    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    c1.fill.solid(); c1.fill.fore_color.rgb = COLOR_CARD; c1.line.color.rgb = COLOR_PURPLE
    tf = c1.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "SELECTED DOMAIN & CONTEXT"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph(); p.text = "\n• Domain: Automated Machine Learning (AutoML) & Enterprise AI Governance.\n• Need: Modern ML deployment requires more than raw predictive power — it demands exploratory data analysis, data quality validation, algorithmic fairness checks, and regulatory compliance.\n• Current Gap: Conventional AutoML frameworks act as opaque 'black boxes', optimizing metrics while completely ignoring exploratory data profiling, fairness auditing, and human governance gates."

    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    c2.fill.solid(); c2.fill.fore_color.rgb = COLOR_CARD; c2.line.color.rgb = COLOR_PURPLE
    tf = c2.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "PROBLEM DEFINITION & SIGNIFICANCE"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    p = tf.add_paragraph(); p.text = "\n• Problem Statement: Unvetted AI deployment propagates historical demographic bias, risks silent data pipeline failures, and causes regulatory non-compliance (e.g., EU AI Act, Fair Credit Reporting Act).\n• Core Governance Constraints: Mandatory enforcement of Disparate Impact (≥0.80) and Demographic Parity Difference (≤0.10) thresholds.\n• Significance: Establishes a transparent, state-checkpointed multi-agent DAG pipeline where human auditors hold final veto authority over deployment."

    # -------------------------------------------------------------
    # SLIDE 3: Parameter 2 — Literature Review (BLANK SLIDE AS REQUESTED)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Parameter 2: Literature / Patent Review & Analysis of Existing Approaches", "CO2 Mapped")

    c_empty = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.6))
    c_empty.fill.solid(); c_empty.fill.fore_color.rgb = COLOR_CARD; c_empty.line.color.rgb = COLOR_MUTED
    tf = c_empty.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "\n\n[ LITERATURE / PATENT REVIEW — RESERVED BLANK SLIDE ]"; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: Parameter 3 — Objectives, Scope & Expected Outcomes
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Parameter 3: Objectives, Scope & Expected Outcomes (2 Marks)", "CO2 & CO4 Mapped")

    col_w = Inches(3.6)
    b1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), col_w, Inches(5.3))
    b1.fill.solid(); b1.fill.fore_color.rgb = COLOR_CARD; b1.line.color.rgb = COLOR_PURPLE
    tf = b1.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "PROJECT OBJECTIVES"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph(); p.text = "\n1. Architect a 6-node multi-agent pipeline (Data Analysis, Planner, Data, Training, Fairness, Governance Gate).\n2. Implement automated quality retry loops and interactive human governance reroute loops.\n3. Persist an unalterable SQLite audit log (`audit_log.db`)."

    b2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), col_w, Inches(5.3))
    b2.fill.solid(); b2.fill.fore_color.rgb = COLOR_CARD; b2.line.color.rgb = COLOR_PURPLE
    tf = b2.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "PROJECT SCOPE"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_AMBER
    p = tf.add_paragraph(); p.text = "\n• Tabular CSV datasets (Classification & Regression).\n• Exploratory Data Profiling & Interactive Chart Visualizations (Chart.js).\n• Ensemble training (RandomForest, XGBoost, Logistic/Ridge).\n• Subgroup fairness auditing & FastAPI web dashboard."

    b3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.5), col_w, Inches(5.3))
    b3.fill.solid(); b3.fill.fore_color.rgb = COLOR_CARD; b3.line.color.rgb = COLOR_PURPLE
    tf = b3.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "EXPECTED OUTCOMES"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    p = tf.add_paragraph(); p.text = "\n• 80% reduction in manual data science pipeline setup.\n• Zero unvetted models deployed to production environments.\n• 100% persistent, queryable audit trail for corporate & regulatory inspection."

    # -------------------------------------------------------------
    # SLIDE 5: Parameter 4 — Proposed Methodology: 6 Pipeline Nodes
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Parameter 4: Proposed Methodology — 6 Pipeline Nodes (3 Marks)", "CO2 & CO4 Mapped")

    card_m = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    card_m.fill.solid(); card_m.fill.fore_color.rgb = COLOR_CARD; card_m.line.color.rgb = COLOR_PURPLE
    tf = card_m.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    
    p = tf.paragraphs[0]; p.text = "MULTI-AGENT TOPOLOGY & NODE RESPONSIBILITIES"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph(); p.text = "1. Data Analysis Agent (EDA Profiler): Computes column statistics, IQR outliers, Pearson feature correlations, target distributions, and chart payloads."
    p = tf.add_paragraph(); p.text = "2. Planner Agent (Groq LLM): Inspects schema metadata, detects data quality issues, pre-computes quality flags, and outputs a structured JSON plan."
    p = tf.add_paragraph(); p.text = "3. Data Agent (Pandas/Scikit-Learn): Executes data cleaning, handles missing value imputation, applies frequency encoding, and standardizes continuous features."
    p = tf.add_paragraph(); p.text = "4. Training Agent (Ensembles): Label-encodes target classes, fits candidate models (RandomForest, XGBoost, Logistic/Ridge), ranks leaderboard metrics, and extracts SHAP importances."
    p = tf.add_paragraph(); p.text = "5. Fairness Agent (Demographic Auditing): Evaluates Disparate Impact (≥0.80) and Demographic Parity Difference (≤0.10) across demographic subgroups."
    p = tf.add_paragraph(); p.text = "6. Governance Gate (Human Interrupt): Pauses execution, presents multi-agent evaluation results to human auditors, receives custom directives, and governs deployment."

    # -------------------------------------------------------------
    # SLIDE 6: Parameter 4 (Cont.) — Governance Gate & Reroute Logic
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Parameter 4 (Cont.): Governance Gate & Reroute Mechanics (3 Marks)", "CO2 & CO4 Mapped")

    c_g1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.3))
    c_g1.fill.solid(); c_g1.fill.fore_color.rgb = COLOR_CARD; c_g1.line.color.rgb = COLOR_GREEN
    tf = c_g1.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "1. APPROVE ROUTE"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    p = tf.add_paragraph(); p.text = "\n• Action: Human auditor approves model performance and fairness report.\n• Outcome: Formally marks the winning model as APPROVED, saves metadata, and records the deployment event in `audit_log.db`."

    c_g2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.3))
    c_g2.fill.solid(); c_g2.fill.fore_color.rgb = COLOR_CARD; c_g2.line.color.rgb = COLOR_AMBER
    tf = c_g2.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "2. REJECT — PLANNER"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_AMBER
    p = tf.add_paragraph(); p.text = "\n• Action: Human auditor rejects due to data quality concerns.\n• Prompt Injection: Accepts custom text directives (e.g. 'Use median imputation for nulls').\n• Reroute Loop: Injects feedback directly into Groq LLM prompt to generate a revised plan."

    c_g3 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.5), Inches(3.6), Inches(5.3))
    c_g3.fill.solid(); c_g3.fill.fore_color.rgb = COLOR_CARD; c_g3.line.color.rgb = COLOR_PURPLE
    tf = c_g3.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]; p.text = "3. REJECT — TRAINING"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_PURPLE
    p = tf.add_paragraph(); p.text = "\n• Action: Human auditor rejects due to model choice or fairness failure.\n• Model Exclusion: Excludes the current winning model (`rejected_models`).\n• Reroute Loop: Retrains and selects the next best candidate algorithm."

    # -------------------------------------------------------------
    # SLIDE 7: Parameter 5 (Part 1) — DEDICATED PIPELINE FLOW GRAPH WITH EXPLICIT ARROWS & DATA ANALYSIS AGENT
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Parameter 5: LangGraph DAG Flow Graph with Explicit Loop Arrows (3 Marks)", "CO2 & CO4 Mapped")

    # Main Outer Container Box
    outer_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(5.6))
    outer_box.fill.solid(); outer_box.fill.fore_color.rgb = COLOR_CARD; outer_box.line.color.rgb = COLOR_PURPLE; outer_box.line.width = Pt(1.5)

    # Start Node (Pill)
    n_start = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(1.48), Inches(6.933), Inches(0.38))
    n_start.fill.solid(); n_start.fill.fore_color.rgb = COLOR_PURPLE; n_start.line.color.rgb = COLOR_ACCENT
    tf = n_start.text_frame; tf.paragraphs[0].text = "START: CSV Dataset Upload + Target Column & Business Objective"; tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow 0
    a0 = slide7.shapes.add_textbox(Inches(6.4), Inches(1.85), Inches(0.5), Inches(0.2))
    a0.text_frame.paragraphs[0].text = "↓"; a0.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a0.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 0: Data Analysis Agent (NEW NODE)
    n_eda = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(2.05), Inches(7.933), Inches(0.46))
    n_eda.fill.solid(); n_eda.fill.fore_color.rgb = RGBColor(20, 35, 55); n_eda.line.color.rgb = COLOR_CYAN; n_eda.line.width = Pt(1.5)
    tf = n_eda.text_frame; p = tf.paragraphs[0]; p.text = "📊 Data Analysis Agent Node (EDA Profiler & Chart Visualizer)"; p.font.size = Pt(10.5); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Column Statistics | IQR Outliers | Pearson Feature Correlations | Target Distribution Charts"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # Arrow 1
    a1 = slide7.shapes.add_textbox(Inches(6.4), Inches(2.50), Inches(0.5), Inches(0.2))
    a1.text_frame.paragraphs[0].text = "↓"; a1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a1.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 1: Planner Agent
    n_plan = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(2.70), Inches(7.933), Inches(0.46))
    n_plan.fill.solid(); n_plan.fill.fore_color.rgb = RGBColor(30, 41, 59); n_plan.line.color.rgb = COLOR_PURPLE; n_plan.line.width = Pt(1.5)
    tf = n_plan.text_frame; p = tf.paragraphs[0]; p.text = "🧠 Planner Agent Node (Groq LLM)"; p.font.size = Pt(10.5); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Schema Analysis | Precomputed Quality Flags | Model Strategy Proposal Generation"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # Arrow 2
    a2 = slide7.shapes.add_textbox(Inches(6.4), Inches(3.15), Inches(0.5), Inches(0.2))
    a2.text_frame.paragraphs[0].text = "↓"; a2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a2.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 2: Data Agent
    n_data = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(3.35), Inches(7.933), Inches(0.46))
    n_data.fill.solid(); n_data.fill.fore_color.rgb = RGBColor(30, 41, 59); n_data.line.color.rgb = COLOR_BLUE; n_data.line.width = Pt(1.5)
    tf = n_data.text_frame; p = tf.paragraphs[0]; p.text = "🧹 Data Agent Node (Pandas / Scikit-Learn)"; p.font.size = Pt(10.5); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Missing Value Imputation | Frequency Encoding | Categorical & Continuous Feature Scaling"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # Arrow 3
    a3 = slide7.shapes.add_textbox(Inches(6.4), Inches(3.80), Inches(0.5), Inches(0.2))
    a3.text_frame.paragraphs[0].text = "↓"; a3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a3.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 3: Training Agent
    n_train = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(4.00), Inches(7.933), Inches(0.46))
    n_train.fill.solid(); n_train.fill.fore_color.rgb = RGBColor(30, 41, 59); n_train.line.color.rgb = COLOR_ACCENT; n_train.line.width = Pt(1.5)
    tf = n_train.text_frame; p = tf.paragraphs[0]; p.text = "🏆 Training Agent Node (Ensemble Framework)"; p.font.size = Pt(10.5); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Target Class LabelEncoding | RandomForest & XGBoost Fitting | SHAP Feature Importance"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # Arrow 4
    a4 = slide7.shapes.add_textbox(Inches(6.4), Inches(4.45), Inches(0.5), Inches(0.2))
    a4.text_frame.paragraphs[0].text = "↓"; a4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a4.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 4: Fairness Agent
    n_fair = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(4.65), Inches(7.933), Inches(0.46))
    n_fair.fill.solid(); n_fair.fill.fore_color.rgb = RGBColor(30, 41, 59); n_fair.line.color.rgb = COLOR_GREEN; n_fair.line.width = Pt(1.5)
    tf = n_fair.text_frame; p = tf.paragraphs[0]; p.text = "⚖️ Fairness Agent Node (Demographic Audit)"; p.font.size = Pt(10.5); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Subgroup Disparate Impact (≥0.80) | Demographic Parity Difference (≤0.10)"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # Arrow 5
    a5 = slide7.shapes.add_textbox(Inches(6.4), Inches(5.10), Inches(0.5), Inches(0.2))
    a5.text_frame.paragraphs[0].text = "↓"; a5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER; a5.text_frame.paragraphs[0].font.size = Pt(10)

    # Node 5: Governance Gate (Diamond/Card Highlight)
    n_gate = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(5.30), Inches(8.733), Inches(0.52))
    n_gate.fill.solid(); n_gate.fill.fore_color.rgb = RGBColor(55, 30, 0); n_gate.line.color.rgb = COLOR_AMBER; n_gate.line.width = Pt(2)
    tf = n_gate.text_frame; p = tf.paragraphs[0]; p.text = "⏸️ Governance Gate Node (Human Interrupt Checkpoint)"; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Pauses execution | Presents Multi-Agent Results | Receives Custom Directives & Reroutes"; p2.font.size = Pt(8); p2.font.color.rgb = COLOR_MUTED; p2.alignment = PP_ALIGN.CENTER

    # --- EXPLICIT LOOP ARROWS & CALLOUT BOXES ---

    # 1. LEFT SIDE: Automated Quality Retry Loop Arrow Box (Data -> Planner)
    loop_left = slide7.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(1.0), Inches(2.80), Inches(1.5), Inches(0.9))
    loop_left.fill.solid(); loop_left.fill.fore_color.rgb = RGBColor(60, 30, 0); loop_left.line.color.rgb = COLOR_AMBER; loop_left.line.width = Pt(1.5)
    tf = loop_left.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.05)
    p = tf.paragraphs[0]; p.text = "⬆ AUTO LOOP\nData ➔ Planner\n(Quality Fail ≤2x)"; p.font.size = Pt(7); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.alignment = PP_ALIGN.CENTER

    # 2. RIGHT SIDE: Human Governance Reroute Loop Arrow Box 1 (Gate -> Planner)
    loop_r1 = slide7.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(10.8), Inches(2.70), Inches(1.4), Inches(1.2))
    loop_r1.fill.solid(); loop_r1.fill.fore_color.rgb = RGBColor(40, 15, 40); loop_r1.line.color.rgb = COLOR_AMBER; loop_r1.line.width = Pt(1.5)
    tf = loop_r1.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.05)
    p = tf.paragraphs[0]; p.text = "⬆ HUMAN LOOP\nGate ➔ Planner\n(Directives ≤2x)"; p.font.size = Pt(7); p.font.bold = True; p.font.color.rgb = COLOR_AMBER; p.alignment = PP_ALIGN.CENTER

    # 3. RIGHT SIDE: Human Governance Reroute Loop Arrow Box 2 (Gate -> Training)
    loop_r2 = slide7.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(10.8), Inches(4.10), Inches(1.4), Inches(1.2))
    loop_r2.fill.solid(); loop_r2.fill.fore_color.rgb = RGBColor(40, 15, 40); loop_r2.line.color.rgb = COLOR_PURPLE; loop_r2.line.width = Pt(1.5)
    tf = loop_r2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.05)
    p = tf.paragraphs[0]; p.text = "⬆ HUMAN LOOP\nGate ➔ Training\n(Exclusion ≤2x)"; p.font.size = Pt(7); p.font.bold = True; p.font.color.rgb = COLOR_PURPLE; p.alignment = PP_ALIGN.CENTER

    # Bottom Decision Legend
    b_leg = slide7.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(11.333), Inches(0.9))
    tf = b_leg.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "DECISION BRANCHES: 🟢 APPROVE ➔ Log Deployment (`audit_log.db`) | 🟠 REJECT DATA ➔ Loop to Planner with Directives | 🔴 REJECT MODEL ➔ Loop to Training with Exclusion"; p.font.size = Pt(8.5); p.font.bold = True; p.font.color.rgb = COLOR_TEXT; p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 8: Parameter 5 (Part 2) — DEDICATED FULL-SLIDE SYSTEM ARCHITECTURE CONTENT
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Parameter 5 (Cont.): System Architecture & Technical Specifications (3 Marks)", "CO2 & CO4 Mapped")

    c_arch1 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    c_arch1.fill.solid(); c_arch1.fill.fore_color.rgb = COLOR_CARD; c_arch1.line.color.rgb = COLOR_PURPLE
    tf = c_arch1.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "CORE BACKEND ENGINE & STORAGE"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph(); p.text = "\n• Stateful LangGraph DAG Engine: Manages sequential execution across all 6 nodes with TypedDict state schemas.\n\n• Data Analysis Module (`data_analysis_agent.py`): Performs exploratory data profiling, outlier detection, and Chart.js payload generation.\n\n• SQLite Checkpointer (`pipeline_state.db`): Persists full state snapshots, enabling crash-recovery and clean interrupt pauses at the Governance Gate.\n\n• Audit Logger (`audit_log.db`): Immutable SQLite event store logging agent outputs, human decisions, and feedback text."

    c_arch2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    c_arch2.fill.solid(); c_arch2.fill.fore_color.rgb = COLOR_CARD; c_arch2.line.color.rgb = COLOR_PURPLE
    tf = c_arch2.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "WEB DASHBOARD & INTERFACING PANELS"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    p = tf.add_paragraph(); p.text = "\n• Enterprise Glassmorphism UI (`static/index.html`): Modern dark slate layout featuring 4 dedicated pages:\n\n  1. Setup & Topology Page: Drag-and-drop CSV uploader, target dropdown, task-type auto-detection, and live DAG topology map.\n\n  2. Data Analysis & Profiling Page: Interactive visual charts (Chart.js), EDA column statistics, feature correlations, and target analysis.\n\n  3. Governance Results Page: Winner KPI card, multi-agent evaluation tabs, and Governance Gate decision panel.\n\n  4. Audit Trail Page: Interactive chronological execution log."

    # -------------------------------------------------------------
    # SLIDE 9: Parameter 6 — Technical Feasibility & Risk Management
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "Parameter 6: Technical Feasibility & Risk Management (3 Marks)", "CO1, CO3 & CO4 Mapped")

    c_risk = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    c_risk.fill.solid(); c_risk.fill.fore_color.rgb = COLOR_CARD; c_risk.line.color.rgb = COLOR_AMBER
    tf = c_risk.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "SYSTEM FEASIBILITY & TECHNICAL GUARDRAILS"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_AMBER
    p = tf.add_paragraph(); p.text = "\n• Performance & Latency: Proven 2.5–4.5 second end-to-end execution latency. Fast parallel model fitting (`n_jobs=-1`) and compressed Groq LLM prompts."
    p = tf.add_paragraph(); p.text = "• Groq 413 Token Overflow Mitigation: Smart schema truncation caps dataset summaries at 40 priority columns (target, sensitive hints, high-nulls), reducing prompt tokens from 10,166 to < 1,800 on 1.4 GB datasets."
    p = tf.add_paragraph(); p.text = "• Target Encoding Guardrail: Automatic `LabelEncoder` conversion ensures 100% training compatibility for XGBoost across binary and multi-class targets."
    p = tf.add_paragraph(); p.text = "• Ethical & Compliance Safeguards: Mandatory subgroup fairness auditing (Disparate Impact & Demographic Parity Difference) and enforced human governance gates."

    # -------------------------------------------------------------
    # SLIDE 10: Parameter 6 (Cont.) — Work Planning & Responsibility Matrix
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "Parameter 6 (Cont.): Work Planning & Module Breakdown (3 Marks)", "CO1, CO3 & CO4 Mapped")

    c_work = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    c_work.fill.solid(); c_work.fill.fore_color.rgb = COLOR_CARD; c_work.line.color.rgb = COLOR_PURPLE
    tf = c_work.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "MODULE BREAKDOWN & RESPONSIBILITY MATRIX"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_ACCENT
    p = tf.add_paragraph(); p.text = "\n• Module 1 (State & Orchestration): LangGraph DAG schema, state serialisation helpers (`df_to_bytes`/`bytes_to_df`), SQLite checkpointer."
    p = tf.add_paragraph(); p.text = "• Module 2 (Agent Intelligence): Data Analysis Agent (EDA/Chart.js), Groq LLM Planner, deterministic Data Agent, Ensemble Trainer with SHAP, Demographic Fairness Agent."
    p = tf.add_paragraph(); p.text = "• Module 3 (Governance & Interrupts): Governance Gate node, custom directive prompt injection, model candidate exclusion logic."
    p = tf.add_paragraph(); p.text = "• Module 4 (REST API & Web UI): FastAPI REST endpoints, SQLite audit logger, 4-page dark slate glassmorphism web dashboard."

    # -------------------------------------------------------------
    # SLIDE 11: Parameter 7 — Preliminary Results & Demonstration
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11)
    add_header(slide11, "Parameter 7: Preliminary Results & System Demonstration (3 Marks)", "CO3 & CO4 Mapped")

    c_demo = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    c_demo.fill.solid(); c_demo.fill.fore_color.rgb = COLOR_CARD; c_demo.line.color.rgb = COLOR_GREEN
    tf = c_demo.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]; p.text = "SYSTEM VERIFICATION & EXPERIMENTAL OUTCOMES"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = COLOR_GREEN
    p = tf.add_paragraph(); p.text = "\n• Empirical Dataset Testing: Validated across Wine Quality (`winetrain-cluster.csv`), Adult Income, and Walmart Sales datasets."
    p = tf.add_paragraph(); p.text = "• Verification Outcomes: Demonstrated multi-model training, dynamic metric leaderboards (Accuracy vs RMSE), SHAP feature importance extraction, and demographic fairness checking."
    p = tf.add_paragraph(); p.text = "• Model Rejection Reroute Verification: Rejecting a winning model excludes it from state (`rejected_models`) and trains remaining candidates on the next pass."
    p = tf.add_paragraph(); p.text = "• Live Application: Fully operational FastAPI REST server (`server.py`) running live at `http://localhost:8000`."

    prs.save("/Users/ruhannp.n/Desktop/mini_project/AegisML_Review2_Presentation.pptx")
    print("Successfully updated AegisML_Review2_Presentation.pptx with Data Analysis Agent and visible loop arrows!")

if __name__ == "__main__":
    create_presentation()
